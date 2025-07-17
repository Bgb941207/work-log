import os
import sys
import zipfile
import re
from pathlib import Path
import win32com.client
from pptx import Presentation
from openpyxl import Workbook

# === 參數區 ===
# 以下參數供日後統一管理
# 輸入
INPUT_FOLDER = None
# 輸出
OUTPUT_ROOT = r"C:\Users\USER\Desktop\ppt_2_excel\八里廠_維修項目整理"
# 圖片格式
IMG_FORMAT = 'jpg'
# 關鍵字範圍
START_KEYWORD = "工作項目說明"
END_KEYWORD   = "工作計畫說明"


def sanitize_text(text):
    """
    移除 Excel 不支援的控制字元 (0x00-0x1F，排除 \t,\n,\r)。
    """
    return ''.join(
        c for c in text
        if c in ('\t', '\n', '\r') or ord(c) >= 32
    )


def find_slide_range(ppt_path, start_kw, end_kw):
    
    try:
        prs = Presentation(ppt_path)
    except (zipfile.BadZipFile, Exception) as e:
        print(f"[Warning] 無法開啟或解析 PPT '{ppt_path}'：{e}")
        return None, None
    start_idx, end_idx = None, None
    for idx, slide in enumerate(prs.slides, start=1):
        text = "".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text
        )
        if start_idx is None and start_kw in text:
            start_idx = idx
        if start_idx is not None and end_kw in text:
            end_idx = idx
            break
    return start_idx, end_idx


def extract_text_to_excel(ppt_path, excel_path, start_slide, end_slide):
    
    try:
        prs = Presentation(ppt_path)
    except (zipfile.BadZipFile, Exception) as e:
        print(f"[Warning] 無法開啟或解析 PPT '{ppt_path}'：{e}")
        wb = Workbook()
        ws = wb.active
        ws.title = "Slides_Text"
        ws.cell(1, 1, "Result")
        ws.cell(2, 1, "無")
        wb.save(excel_path)
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Slides_Text"
    ws.cell(1, 1, "Slide Number")
    ws.cell(1, 2, "Text Content")
    row = 2
    for idx, slide in enumerate(prs.slides, start=1):
        if idx < start_slide or idx > end_slide:
            continue
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                content = shape.text.strip()
                if content:
                    texts.append(content)
        combined = sanitize_text("\n".join(texts))
        if combined:
            ws.cell(row, 1, idx)
            ws.cell(row, 2, combined)
            row += 1
    if row == 2:
        ws.cell(2, 1, "無")
    wb.save(excel_path)
    print(f"✔ 已儲存文字 Excel: {excel_path}")


def export_images(ppt_path, output_dir, img_format, start_slide, end_slide):
   
    ppt_path = os.path.abspath(ppt_path)
    if not os.path.exists(ppt_path):
        print(f"[Error] 檔案不存在: {ppt_path}", file=sys.stderr)
        return
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(ppt_path, ReadOnly=1, WithWindow=0)
    except Exception as e:
        print(f"[Warning] COM 開啟 PPT 失敗 '{ppt_path}'：{e}")
        return
    os.makedirs(output_dir, exist_ok=True)
    total = pres.Slides.Count
    end_idx = min(end_slide, total)
    for idx in range(start_slide, end_idx + 1):
        slide = pres.Slides.Item(idx)
        out_path = os.path.join(output_dir, f"slide_{idx}.{img_format}")
        slide.Export(out_path, img_format)
        print(f"✔ 匯出圖片 第{idx}頁 → {out_path}")
    pres.Close()
    app.Quit()


if __name__ == '__main__':
    # 互動輸入：僅需填入 INPUT_FOLDER
    INPUT_FOLDER = input("請輸入包含 PPT 檔案的資料夾絕對路徑：").strip().strip('"\'')
    folder = Path(INPUT_FOLDER)
    if not folder.is_dir():
        print(f"[Error] 資料夾不存在: {INPUT_FOLDER}")
        sys.exit(1)

    # 遞迴處理所有 .ppt/.pptx 檔案
    for ppt_file in folder.rglob('*.ppt*'):
        orig_name = ppt_file.stem
        m = re.match(r'^(\d+)', orig_name)
        base_name = m.group(1) if m else orig_name
        out_dir = Path(OUTPUT_ROOT) / base_name
        out_dir.mkdir(parents=True, exist_ok=True)
        excel_path = out_dir / f"{base_name}.xlsx"

        # 偵測範圍
        start_idx, end_idx = find_slide_range(str(ppt_file), START_KEYWORD, END_KEYWORD)
        if start_idx is None or end_idx is None:
            wb = Workbook()
            ws = wb.active
            ws.title = "Slides_Text"
            ws.cell(1, 1, "Result")
            ws.cell(2, 1, "無")
            wb.save(excel_path)
            print(f"⚠ [{base_name}] 找不到或無法解析 PPT，僅輸出 Excel '無'。")
        else:
            print(f"[{base_name}] 偵測到範圍：第{start_idx}頁 → 第{end_idx}頁")
            extract_text_to_excel(str(ppt_file), str(excel_path), start_idx, end_idx)
            export_images(str(ppt_file), str(out_dir), IMG_FORMAT, start_idx, end_idx)

    print("\n全部檔案處理完成，輸出位於：", OUTPUT_ROOT)

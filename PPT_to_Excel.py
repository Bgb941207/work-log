import os
import sys
import zipfile
import re
from pathlib import Path
from win32com.client import gencache
from pptx import Presentation
from openpyxl import Workbook

# === 參數區 ===
# 互動輸入的 PPT 資料夾
INPUT_FOLDER = None
# 輸出根目錄
OUTPUT_ROOT = r"  "
# 圖片格式 (jpg, png, etc.)
IMG_FORMAT = 'jpg'
# 關鍵字範圍
START_KEYWORD = "工作項目說明" #本月工作項目報告
END_KEYWORD   = "工作計畫說明" #下個月工作計劃說明


def sanitize_text(text):
    """移除 Excel 不支援的控制字元 (0x00-0x1F，排除 \t,\n,\r)。"""
    return ''.join(
        c for c in text
        if c in ('\t', '\n', '\r') or ord(c) >= 32
    )


def find_slide_range(ppt_path, start_kw, end_kw):
    """找出 PPT 中包含 start_kw 與 end_kw 的第一個範圍 (起訖頁)。"""
    try:
        prs = Presentation(ppt_path)
    except (zipfile.BadZipFile, Exception) as e:
        print(f"[Warning] 無法開啟或解析 PPT '{ppt_path}'：{e}")
        return None, None

    start_idx = end_idx = None
    for idx, slide in enumerate(prs.slides, start=1):
        text = "".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text
        )
        if start_idx is None and start_kw in text:
            start_idx = idx
        if start_idx is not None and end_kw in text:
            end_idx = idx
            break
    return start_idx, end_idx


def extract_text_to_excel(ppt_path, excel_path, start_slide, end_slide, image_dir=None, img_format='jpg'):
    """
    將指定頁碼範圍內的文字匯出到 Excel，
    並於 C 欄填入每張截圖的絕對路徑 (若提供 image_dir)。
    """
    try:
        prs = Presentation(ppt_path)
    except (zipfile.BadZipFile, Exception) as e:
        print(f"[Warning] 無法開啟或解析 PPT '{ppt_path}'：{e}")
        wb = Workbook()
        ws = wb.active
        ws.title = "Slides_Text"
        ws.cell(1, 1, "Result")
        ws.cell(2, 1, "無")
        wb.save(excel_path)
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Slides_Text"
    # 標頭
    ws.cell(1, 1, "Slide Number")
    ws.cell(1, 2, "Text Content")
    ws.cell(1, 3, "Image Path")

    row = 2
    for idx, slide in enumerate(prs.slides, start=1):
        if idx < start_slide or idx > end_slide:
            continue

        # 收集文字
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                content = shape.text.strip()
                if content:
                    texts.append(content)
        combined = sanitize_text("\n".join(texts)) or "無文字"

        # 填入 A、B
        ws.cell(row, 1, idx)
        ws.cell(row, 2, combined)

        # 如有提供 image_dir，計算截圖檔名並寫入 C 欄
        if image_dir:
            img_name = f"slide_{idx}.{img_format}"
            img_path = os.path.abspath(os.path.join(image_dir, img_name))
            ws.cell(row, 3, img_path)

        row += 1

    wb.save(excel_path)
    print(f"✔ 已儲存文字與截圖路徑至 Excel: {excel_path}")


def export_images(ppt_path, output_dir, img_format, start_slide, end_slide):
    """使用 COM 介面匯出 PPT 範圍內的每頁為圖片。"""
    ppt_path = os.path.abspath(ppt_path)
    if not os.path.exists(ppt_path):
        print(f"[Error] 檔案不存在: {ppt_path}", file=sys.stderr)
        return

    # 早期綁定 PowerPoint
    app = gencache.EnsureDispatch("PowerPoint.Application")
    # WithWindow=0 已經隱藏視窗，不需設定 Visible
    pres = app.Presentations.Open(ppt_path, ReadOnly=1, WithWindow=0)

    os.makedirs(output_dir, exist_ok=True)
    total = pres.Slides.Count
    end_idx = min(end_slide, total)

    for idx in range(start_slide, end_idx + 1):
        slide = pres.Slides(idx)
        out_path = os.path.join(output_dir, f"slide_{idx}.{img_format}")
        slide.Export(out_path, img_format)
        print(f"✔ 匯出圖片 第{idx}頁 → {out_path}")

    pres.Close()
    app.Quit()


if __name__ == '__main__':
    # 互動輸入：僅需填入 INPUT_FOLDER
    INPUT_FOLDER = input("請輸入包含 PPT 檔案的資料夾絕對路徑：").strip().strip('"\'')
    folder = Path(INPUT_FOLDER)
    if not folder.is_dir():
        print(f"[Error] 資料夾不存在: {INPUT_FOLDER}")
        sys.exit(1)

    # 遞迴處理所有 .ppt/.pptx 檔案
    for ppt_file in folder.rglob('*.ppt*'):
        orig_name = ppt_file.stem
        m = re.match(r'^(\d+)', orig_name)
        base_name = m.group(1) if m else orig_name

        out_dir = Path(OUTPUT_ROOT) / base_name
        out_dir.mkdir(parents=True, exist_ok=True)

        excel_path = out_dir / f"{base_name}.xlsx"

        # 偵測範圍
        start_idx, end_idx = find_slide_range(str(ppt_file), START_KEYWORD, END_KEYWORD)
        if start_idx is None or end_idx is None:
            wb = Workbook()
            ws = wb.active
            ws.title = "Slides_Text"
            ws.cell(1, 1, "Result")
            ws.cell(2, 1, "無")
            wb.save(excel_path)
            print(f"⚠ [{base_name}] 找不到範圍，僅輸出 Excel '無'。")
        else:
            print(f"[{base_name}] 偵測到範圍：第{start_idx}頁 → 第{end_idx}頁")
            extract_text_to_excel(
                str(ppt_file), str(excel_path),
                start_idx, end_idx,
                image_dir=str(out_dir), img_format=IMG_FORMAT
            )
            export_images(
                str(ppt_file), str(out_dir),
                IMG_FORMAT, start_idx, end_idx
            )

    print("\n全部檔案處理完成，輸出位於：", OUTPUT_ROOT)
    
